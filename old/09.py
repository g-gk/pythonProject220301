import collections
import collections.abc
import random
from pptx import Presentation

# from pptx.util import Inches
# from pptx.util import Pt
# import os


# создаем новую презентацию
prs = Presentation()
# получаем схему расположения элементов для заголовочного слайда
title_slide_layout = prs.slide_layouts[0]
title_slide_layout2 = prs.slide_layouts[1]
# создаем заголовочный слайд
slide = prs.slides.add_slide(title_slide_layout)
slide2 = prs.slides.add_slide(title_slide_layout2)
slide3 = prs.slides.add_slide(title_slide_layout2)
slide4 = prs.slides.add_slide(title_slide_layout2)
slide5 = prs.slides.add_slide(title_slide_layout2)
slide6 = prs.slides.add_slide(title_slide_layout2)
# создаем у слайда заголовок и текст
title = slide.shapes.title
subtitle = slide.placeholders[1]
title2 = slide2.shapes.title
subtitle2 = slide2.placeholders[1]
title3 = slide3.shapes.title
subtitle3 = slide3.placeholders[1]
title4 = slide4.shapes.title
subtitle4 = slide4.placeholders[1]
title5 = slide5.shapes.title
subtitle5 = slide5.placeholders[1]
title6 = slide6.shapes.title
subtitle6 = slide6.placeholders[1]

title.text = "Привет"
subtitle.text = "Это презентация про модуль random. Немного про методы модуля:"
title2.text = random.random.__name__
subtitle2.text = random.random.__doc__
title3.text = random.randint.__name__
subtitle3.text = random.randint.__doc__
title4.text = random.choice.__name__
subtitle4.text = random.choice.__doc__
title5.text = random.choices.__name__
subtitle5.text = random.choices.__doc__
title6.text = random.shuffle.__name__
subtitle6.text = random.shuffle.__doc__
# сохраняем презентацию
prs.save('test.pptx')
# os.startfile('C:\\Users\\chane\\Desktop\\homework\\test.pptx')
print('Пытался изменить шрифт, но ничего не вышло; '
      'брал код с документации, но выдавало ошибки')
