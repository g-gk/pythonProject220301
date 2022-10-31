import collections
import collections.abc

from pptx import Presentation
from pptx.util import Pt

from random import *
import pydoc
import string


def del_bs(s):
    res = []
    for c in s:
        if c == '\x08':
            res = res[:-1]
        else:
            res.append(c)
    return "".join(res)


# создаем новую презентацию
prs = Presentation()
# получаем схему расположения элементов для заголовочного слайда
title_slide_layout = prs.slide_layouts[1]
# создаем заголовочный слайд
slide = prs.slides.add_slide(title_slide_layout)
# получем описание модуля
help_text = Random.__doc__
# создаем у слайда заголовок и текст
title = slide.shapes.title
p = title.text_frame.paragraphs[0]
run = p.add_run()
run.text = "Random"
font = run.font
font.name = 'Courier New'
font.size = Pt(14)
font.bold = True

subtitle = slide.placeholders[1]
p = subtitle.text_frame.paragraphs[0]
run = p.add_run()
run.text = help_text
font = run.font
font.name = 'Courier New'
font.size = Pt(12)

function_list = [method for method in dir(
    Random) if not method.startswith('_') and method != method.upper()]
slide_count = 0

for function in function_list:
    if slide_count >= 5:  # если число слайдов достаточное
        break  # заканчиваем перебор методов

    help_text = del_bs(pydoc.render_doc(getattr(Random, function), "%s"))

    # формируем слайд по функции
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    p = title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = function
    font = run.font
    font.name = 'Courier New'
    font.size = Pt(14)
    font.bold = True

    subtitle = slide.placeholders[1]

    p = subtitle.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = help_text
    font = run.font
    font.name = 'Courier New'
    font.size = Pt(12)

    slide_count += 1

prs.save('res.pptx')
