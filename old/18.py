from pptx import Presentation

# создаем презентацию
prs = Presentation()

# первый слайд
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "choice()"
subtitle.text = 'random.choice = choice(seq) method of random.Random instance.' \
                'Choose a random element from a non-empty sequence.'

# второй слайд
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "randint()"
subtitle.text = 'random.randint = randint(a, b) method of random.Random instance.' \
                'Return random integer in range [a, b], including both end points.'

# третий слайд
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "setstate()"
subtitle.text = 'random.setstate = setstate(state) method of random.Random instance.' \
                'Restore internal state from object returned by getstate().'

# четвертый слайд
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "uniform()"
subtitle.text = 'random.uniform = uniform(a, b) method of random.Random instance.' \
                'Get a random number in the range [a, b) or [a, b] depending on rounding.'

# пятый слайд
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "paretovariate()"
subtitle.text = 'random.paretovariate = paretovariate(alpha) method of random.Random instance.' \
                'Pareto distribution.  alpha is the shape parameter.'

# сохраняем презентацию
prs.save('res.pptx')
